#!/usr/bin/env python3
# ─────────────────────────────────────────────────────────────────────────────
# SlideFluid conversion backend — the entire Python side, one file.
#
# ⚠ Read the right topic doc before editing:
#     • PDF → PPTX path .................. docs/PDF_FILL.md
#     • DOCX / TXT / RTF → PPTX path ..... docs/DOCX_TXT.md
#     • PPTX → notes / split / companion . docs/PPTX_PIPELINE.md
#     • IPC schema + CLI flag contract ... docs/IPC.md
#
# This module is invoked as a subprocess by `src/main.js` (`ConversionJob`).
# In `--ipc` mode it emits newline-delimited JSON to stdout, one event per
# line; nothing else may go to stdout in that mode (a stray `print` will
# break parsing on the renderer side — `_emit` is the ONLY exit channel).
#
# Key invariants:
#   • All file writes are atomic — `_atomic_save_pptx` / `_atomic_write_text`.
#     Never write directly to the final path; PowerPoint opening a half-saved
#     file is a real failure mode we already paid for once.
#   • The CLI flag names (`--fill`, `--slide-theme`, `--text-align`,
#     `--split-choices`, `--pptx-mode`, `--pptx-threshold`, etc.) are the IPC
#     contract. Renaming any of them requires updating
#     `ConversionJob._buildArgs` in main.js AND docs/IPC.md.
#   • python-pptx 1.0 has surprises (Part.load arg order, get_or_add return
#     type, external rel handling). See docs/PPTX_PIPELINE.md before assuming
#     any signature.
#   • Don't `print()` in `--ipc` mode. Use `_emit` only.
# ─────────────────────────────────────────────────────────────────────────────
"""SlideFluid conversion backend.

Invoke modes:
    Conversion (batch):    slidefluid_convert.py [flags] file1 [file2 ...]
    Probe — docx info:     slidefluid_convert.py --docx-info <file>
    Probe — pptx info:     slidefluid_convert.py --pptx-info <file>
    Probe — docx analyze:  slidefluid_convert.py --analyze <file>
    Preflight check:       slidefluid_convert.py --preflight

Full CLI flag inventory and IPC event schema: docs/IPC.md.
"""

import argparse
import copy
import json
import math
import os
import re
import sys
import tempfile
from pathlib import Path

import numpy as np
from PIL import Image, ImageFilter
from pdf2image import convert_from_path
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5
SLIDE_WIDTH_EMU = int(SLIDE_WIDTH_IN * 914400)   # 12,192,120
SLIDE_HEIGHT_EMU = int(SLIDE_HEIGHT_IN * 914400)  # 6,858,000

TARGET_AR = SLIDE_WIDTH_IN / SLIDE_HEIGHT_IN  # 1.7773…

# ---------------------------------------------------------------------------
# IPC / logging helpers
# ---------------------------------------------------------------------------

_ipc_mode = False


def _emit(obj: dict):
    """Write a JSON line to stdout (IPC mode) or human-readable text (CLI mode)."""
    if _ipc_mode:
        print(json.dumps(obj), flush=True)
    else:
        t = obj.get("type", "")
        if t == "start":
            print(f"\n[{obj['file_index']}/{obj['total_files']}] {obj['file']}")
        elif t == "progress":
            print(f"  page {obj['page']}/{obj['total_pages']}  {obj['message']}", end="\r")
        elif t == "done":
            print(f"\n  ✓ Done — {obj['slides']} slides → {obj['output']}")
        elif t == "error":
            print(f"\n  ✗ Error: {obj['message']}", file=sys.stderr)
        elif t == "batch_done":
            print(
                f"\nBatch complete: {obj['converted']} converted, "
                f"{obj['skipped']} skipped, {obj['errors']} errors, "
                f"{obj['total_slides']} total slides."
            )
        elif t == "warn":
            print(f"  ! {obj['message']}")


# ---------------------------------------------------------------------------
# Aspect-ratio helpers
# ---------------------------------------------------------------------------

def detect_ar(width_px: int, height_px: int) -> str:
    """Return a human-readable aspect ratio tag."""
    ar = width_px / height_px
    if abs(ar - TARGET_AR) < 0.02:
        return "16:9"
    if abs(ar - (4 / 3)) < 0.02:
        return "4:3"
    return f"{width_px}:{height_px}"


def is_native_169(width_px: int, height_px: int) -> bool:
    ar = width_px / height_px
    return abs(ar - TARGET_AR) < 0.02


# ---------------------------------------------------------------------------
# Pillarbox fill modes
# ---------------------------------------------------------------------------

def _bar_dims(img: Image.Image) -> tuple[int, int, int]:
    """
    For a non-16:9 image that will be fit-to-height on a 16:9 canvas,
    return (canvas_width, img_width_scaled, bar_width_each).
    The image is scaled to canvas height; bars flank it horizontally.
    """
    src_ar = img.width / img.height
    canvas_h = img.height  # we keep height, vary width
    canvas_w = int(canvas_h * TARGET_AR)
    img_w_scaled = int(canvas_h * src_ar)
    bar_w = (canvas_w - img_w_scaled) // 2
    return canvas_w, img_w_scaled, bar_w


def fill_black(img: Image.Image) -> Image.Image:
    """Pillarbox with solid black bars."""
    canvas_w, img_w, bar_w = _bar_dims(img)
    canvas = Image.new("RGB", (canvas_w, img.height), (0, 0, 0))
    canvas.paste(img.resize((img_w, img.height), Image.LANCZOS), (bar_w, 0))
    return canvas


def fill_color_match(img: Image.Image) -> Image.Image:
    """
    Sample the left and right edge columns of the image, average them,
    and fill the pillarbox bars with that flat color.
    Falls back to black on any error.
    """
    try:
        arr = np.array(img)
        # Sample outermost 5px on each side
        sample_w = min(5, img.width // 4)
        left_strip = arr[:, :sample_w, :3]
        right_strip = arr[:, -sample_w:, :3]
        combined = np.concatenate([left_strip.reshape(-1, 3),
                                   right_strip.reshape(-1, 3)], axis=0)
        avg = tuple(int(x) for x in combined.mean(axis=0))

        canvas_w, img_w, bar_w = _bar_dims(img)
        canvas = Image.new("RGB", (canvas_w, img.height), avg)
        canvas.paste(img.resize((img_w, img.height), Image.LANCZOS), (bar_w, 0))
        return canvas
    except Exception:
        return fill_black(img)


def fill_smear(img: Image.Image) -> Image.Image:
    """
    Blur-extend the outermost columns outward into the pillarbox bars.
    Creates a natural-looking smear effect. Falls back to color_match on error.
    """
    try:
        canvas_w, img_w, bar_w = _bar_dims(img)
        if bar_w <= 0:
            return img

        img_scaled = img.resize((img_w, img.height), Image.LANCZOS)
        arr = np.array(img_scaled)

        # Build left bar: repeat left edge column, then blur
        left_col = arr[:, :1, :]                            # (H, 1, 3)
        left_bar_arr = np.repeat(left_col, bar_w, axis=1)   # (H, bar_w, 3)

        # Build right bar: repeat right edge column, then blur
        right_col = arr[:, -1:, :]
        right_bar_arr = np.repeat(right_col, bar_w, axis=1)

        # Assemble full canvas array
        canvas_arr = np.concatenate([left_bar_arr, arr, right_bar_arr], axis=1)
        canvas = Image.fromarray(canvas_arr.astype(np.uint8), "RGB")

        # Apply a strong horizontal blur to the bar regions only
        blurred = canvas.filter(ImageFilter.GaussianBlur(radius=bar_w // 3 + 4))

        # Composite: use blurred for bars, original scaled img for center
        result = blurred.copy()
        result.paste(img_scaled, (bar_w, 0))
        return result
    except Exception:
        return fill_color_match(img)


FILL_FUNCS = {
    "black": fill_black,
    "color_match": fill_color_match,
    "smear": fill_smear,
}


# ---------------------------------------------------------------------------
# Core page → slide
# ---------------------------------------------------------------------------

def add_page_to_pptx(
    prs: Presentation,
    page_img: Image.Image,
    fill_mode: str,
) -> None:
    """Rasterize one PDF page, apply pillarbox if needed, add to PPTX."""
    w, h = page_img.size
    native = is_native_169(w, h)

    if native:
        final_img = page_img
    else:
        final_img = FILL_FUNCS.get(fill_mode, fill_black)(page_img)

    # Write to temp file, add to slide, delete immediately
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        final_img.save(tmp_path, "PNG", optimize=False)

        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        pic = slide.shapes.add_picture(
            tmp_path,
            left=Emu(0),
            top=Emu(0),
            width=Emu(SLIDE_WIDTH_EMU),
            height=Emu(SLIDE_HEIGHT_EMU),
        )
    finally:
        os.unlink(tmp_path)  # delete immediately — never accumulate


# ---------------------------------------------------------------------------
# PDF → PPTX
# ---------------------------------------------------------------------------

def convert_pdf(
    pdf_path: Path,
    output_dir: Path,
    dpi: int = 72,
    fill_mode: str = "black",
    overwrite: bool = False,
    skip_existing: bool = False,
    suffix: str = "",
    file_index: int = 1,
    total_files: int = 1,
    poppler_path: str | None = None,
) -> dict:
    """
    Convert a single PDF to PPTX.

    Returns:
        {"ok": True,  "output": str, "slides": int}  on success
        {"ok": False, "message": str}                 on error
    """
    stem = pdf_path.stem + suffix
    out_path = output_dir / f"{stem}.pptx"

    # --- Overwrite check ---
    if out_path.exists():
        if skip_existing:
            _emit({"type": "warn", "file": str(pdf_path),
                   "message": f"Skipped — {out_path.name} already exists."})
            return {"ok": False, "message": "skipped — file exists", "skipped": True}
        if not overwrite:
            # In CLI mode ask; in IPC mode always overwrite (caller handles this)
            if not _ipc_mode:
                ans = input(f"  {out_path.name} already exists. Overwrite? [y/N] ").strip().lower()
                if ans != "y":
                    return {"ok": False, "message": "skipped — user declined overwrite",
                            "skipped": True}

    _emit({
        "type": "start",
        "file": str(pdf_path),
        "file_index": file_index,
        "total_files": total_files,
    })

    # --- Rasterize ---
    try:
        kwargs = dict(dpi=dpi, fmt="RGB", thread_count=2, use_cropbox=False)
        if poppler_path:
            kwargs["poppler_path"] = poppler_path
        pages = convert_from_path(str(pdf_path), **kwargs)
    except PDFInfoNotInstalledError:
        msg = "Poppler not found — cannot rasterize."
        _emit({"type": "error", "file": str(pdf_path), "message": msg})
        return {"ok": False, "message": msg}
    except PDFPageCountError:
        msg = "Password-protected or unreadable — skipped."
        _emit({"type": "error", "file": str(pdf_path), "message": msg})
        return {"ok": False, "message": msg}
    except PDFSyntaxError as e:
        msg = f"Corrupt PDF — {e}"
        _emit({"type": "error", "file": str(pdf_path), "message": msg})
        return {"ok": False, "message": msg}
    except Exception as e:
        msg = f"Unexpected error during rasterization: {e}"
        _emit({"type": "error", "file": str(pdf_path), "message": msg})
        return {"ok": False, "message": msg}

    total_pages = len(pages)

    # --- Build PPTX ---
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    for i, page_img in enumerate(pages, start=1):
        _emit({
            "type": "progress",
            "file": str(pdf_path),
            "page": i,
            "total_pages": total_pages,
            "message": f"Converting page {i} of {total_pages}",
        })
        try:
            add_page_to_pptx(prs, page_img, fill_mode)
        except Exception as e:
            msg = f"Failed on page {i}: {e}"
            _emit({"type": "error", "file": str(pdf_path), "message": msg})
            return {"ok": False, "message": msg}

    # --- Save ---
    try:
        prs.save(str(out_path))
    except PermissionError:
        msg = f"Cannot write to {out_path} — permission denied."
        _emit({"type": "error", "file": str(pdf_path), "message": msg})
        return {"ok": False, "message": msg}
    except OSError as e:
        msg = f"Disk error saving {out_path.name}: {e}"
        _emit({"type": "error", "file": str(pdf_path), "message": msg})
        return {"ok": False, "message": msg}

    _emit({
        "type": "done",
        "file": str(pdf_path),
        "output": str(out_path),
        "slides": total_pages,
    })
    return {"ok": True, "output": str(out_path), "slides": total_pages}


# ---------------------------------------------------------------------------
# Batch runner
# ---------------------------------------------------------------------------

def run_batch(
    pdf_paths: list[Path],
    output_dir: Path,
    dpi: int = 72,
    fill_mode: str = "black",
    overwrite: bool = False,
    skip_existing: bool = False,
    suffix: str = "",
    poppler_path: str | None = None,
    slide_theme: str = "light",
    text_align: str = "left",
    split_choices: list | None = None,
    pptx_opts: dict | None = None,
) -> int:
    """Run conversion on a list of PDF paths. Returns exit code (0 = all ok)."""
    converted = 0
    skipped = 0
    errors = 0
    total_slides = 0
    total = len(pdf_paths)
    pptx_opts = pptx_opts or {}

    for i, pdf in enumerate(pdf_paths, start=1):
        ext = pdf.suffix.lower()
        if ext == ".pptx":
            result = convert_pptx_notes(
                file_path=pdf,
                output_dir=output_dir,
                mode=pptx_opts.get("mode", "teleprompter"),
                threshold=pptx_opts.get("threshold", 12),
                include_headers=pptx_opts.get("include_headers", True),
                include_slide_numbers=pptx_opts.get("include_slide_numbers", False),
                slide_theme=pptx_opts.get("slide_theme", "dark"),
                suffix=suffix,
                file_index=i,
                total_files=total,
            )
        elif ext in (".txt", ".docx"):
            result = convert_text_doc(
                file_path=pdf,
                output_dir=output_dir,
                suffix=suffix,
                file_index=i,
                total_files=total,
                slide_theme=slide_theme,
                text_align=text_align,
                split_choices=split_choices,
            )
        else:
            result = convert_pdf(
                pdf_path=pdf,
                output_dir=output_dir,
                dpi=dpi,
                fill_mode=fill_mode,
                overwrite=overwrite,
                skip_existing=skip_existing,
                suffix=suffix,
                file_index=i,
                total_files=total,
                poppler_path=poppler_path,
            )
        if result.get("ok"):
            converted += 1
            total_slides += result.get("slides", 0)
        elif result.get("skipped"):
            skipped += 1
        else:
            errors += 1

    _emit({
        "type": "batch_done",
        "converted": converted,
        "skipped": skipped,
        "errors": errors,
        "total_slides": total_slides,
    })
    return 0 if errors == 0 else 1


# ---------------------------------------------------------------------------
# PDF discovery
# ---------------------------------------------------------------------------

_SUPPORTED_EXTS = {".pdf", ".docx", ".txt", ".pptx"}


def collect_files(inputs: list[str]) -> list[Path]:
    """Expand a mix of file paths and folder paths to a flat list of supported files."""
    result = []
    for item in inputs:
        p = Path(item)
        if p.is_dir():
            for ext in (".pdf", ".PDF", ".docx", ".DOCX", ".txt", ".TXT", ".pptx", ".PPTX"):
                result.extend(sorted(p.rglob(f"*{ext}")))
        elif p.is_file():
            if p.suffix.lower() in _SUPPORTED_EXTS:
                result.append(p)
            else:
                print(f"Warning: {p} — unsupported type, skipped.", file=sys.stderr)
        else:
            print(f"Warning: {p} not found — skipped.", file=sys.stderr)
    seen: set[Path] = set()
    deduped = []
    for p in result:
        if p not in seen:
            seen.add(p)
            deduped.append(p)
    return deduped


# Keep old name as alias for any external callers
collect_pdfs = collect_files


# ---------------------------------------------------------------------------
# DOCX / TXT → PPTX
# ---------------------------------------------------------------------------

# Text-fitting constants (in points, 1 pt = 1/72 inch)
_TXT_MARGIN_IN      = 0.55   # left / right margin
_TXT_TOP_MARGIN_IN  = 0.15   # top margin (smaller = more vertical room)
_TXT_W_PT       = (SLIDE_WIDTH_IN  - _TXT_MARGIN_IN     * 2) * 72   # ≈ 881 pt
_TXT_H_PT       = (SLIDE_HEIGHT_IN - _TXT_TOP_MARGIN_IN * 2) * 72   # ≈ 518 pt
_CHAR_W_RATIO   = 0.44   # avg char width as fraction of point size (calibri/system-ui)
_LINE_H_RATIO   = 1.20   # line height as multiple of font size
_MIN_FONT_PT    = 20
_MAX_FONT_PT    = 54


# --- Parsers ---

def _strip_rtf(text: str) -> str:
    """Minimal RTF→plain text stripper. Handles TextEdit/macOS RTF saved with
    a .txt or .rtf extension. Removes control words and groups; preserves
    visible text. Not a full RTF parser — good enough for note-taking use."""
    # Decode \uNNNN unicode escapes first (RTF stores Unicode as decimal)
    def _u(m):
        try:
            return chr(int(m.group(1)) & 0xFFFF)
        except Exception:
            return ""
    text = re.sub(r"\\u(-?\d+)\??", _u, text)
    # Decode \'XX hex escapes (single byte in document codepage, usually CP1252)
    def _h(m):
        try:
            return bytes([int(m.group(1), 16)]).decode('cp1252', errors='replace')
        except Exception:
            return ""
    text = re.sub(r"\\'([0-9a-fA-F]{2})", _h, text)
    # Drop known non-content groups (fonttbl, colortbl, etc.) by brace-depth tracking
    NON_CONTENT_GROUPS = ('fonttbl', 'colortbl', 'expandedcolortbl', 'stylesheet',
                          'pict', 'info', 'header', 'footer', 'listtable', 'listoverridetable',
                          'rsidtbl', 'generator', 'cocoascreenfonts', 'cocoartf')
    for group in NON_CONTENT_GROUPS:
        pattern = re.compile(r'\{\\\*?\\?' + group + r'\b')
        while True:
            m = pattern.search(text)
            if not m:
                break
            start = m.start()
            depth = 0
            end = start
            for i in range(start, len(text)):
                if text[i] == '{':
                    depth += 1
                elif text[i] == '}':
                    depth -= 1
                    if depth == 0:
                        end = i + 1
                        break
            else:
                end = len(text)
            text = text[:start] + text[end:]
    # Drop \par/\line/\tab as actual whitespace
    text = re.sub(r"\\par[d]?\b", "\n", text)
    text = re.sub(r"\\line\b", "\n", text)
    text = re.sub(r"\\tab\b", "\t", text)
    # Drop \* (ignorable destination marker)
    text = re.sub(r"\\\*", "", text)
    # Strip remaining control words (\word, \word123) with optional trailing space
    text = re.sub(r"\\[A-Za-z]+-?\d* ?", "", text)
    # Strip leftover braces (the outer document wrapper, etc.)
    text = text.replace("{", "").replace("}", "")
    # Unescape backslash-escaped chars
    text = re.sub(r"\\([\\'\"])", r"\1", text)
    # Normalise non-breaking spaces and line separators
    text = text.replace("\xa0", " ")
    text = text.replace("\u2028", "\n").replace("\u2029", "\n")
    # Promote inline bullet markers (•, ·, etc.) onto their own lines so each
    # bullet becomes its own paragraph in the slide
    text = re.sub(r"[ \t]+([•·▪►])[ \t]+", r"\n\1 ", text)
    # Collapse 3+ consecutive blank lines to exactly 2 (the slide-break threshold)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text


def _merge_orphan_headings(slides: list) -> list:
    """If a slide is only a single heading paragraph, merge it into the next
    slide. This catches source docs with extra blank lines between heading and
    body content (common in hand-formatted notes)."""
    if not slides:
        return slides
    result = []
    i = 0
    while i < len(slides):
        cur = slides[i]
        non_empty = [p for p in cur if p.get("text", "").strip()]
        if (len(non_empty) == 1 and non_empty[0].get("is_heading")
                and i + 1 < len(slides)):
            # Prepend this heading to the next slide
            result.append(cur + slides[i + 1])
            i += 2
        else:
            result.append(cur)
            i += 1
    return result


def _parse_txt(path: Path) -> tuple[list[list[dict]], list[str]]:
    """
    Split a plain-text file into slides.
    Two or more consecutive blank lines = slide boundary.
    A single blank line is treated as a paragraph break within the same slide.
    """
    text = path.read_text(encoding="utf-8", errors="replace")
    # Detect RTF (TextEdit on macOS saves .rtf even with .txt extension sometimes)
    if text.lstrip().startswith(r"{\rtf"):
        text = _strip_rtf(text)
    # Normalise line endings and non-breaking spaces (common from Word/Google Docs exports)
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\xa0", " ")
    # Split on 2+ consecutive blank lines (lines containing only whitespace count as blank)
    raw_blocks = re.split(r"\n(?:[ \t]*\n){2,}", text.strip())
    slides = []
    BULLET_CHARS = ('•', '·', '▪', '►', '-', '*')
    for block in raw_blocks:
        block = block.strip()
        if not block:
            continue
        paragraphs = []
        found_heading = False
        for line in block.split("\n"):
            line = line.rstrip()
            if not line:
                continue
            is_bullet = line.lstrip().startswith(BULLET_CHARS)
            if is_bullet:
                # Strip the leading bullet character + spaces — the slide
                # builder adds its own "• " prefix
                line = line.lstrip().lstrip(''.join(BULLET_CHARS)).strip()
                if not line:
                    continue
            # First non-bullet line of the slide is treated as the heading
            is_heading = (not is_bullet) and (not found_heading)
            if is_heading:
                found_heading = True
            paragraphs.append({
                "text": line,
                "runs": [{"text": line, "bold": False, "italic": False}],
                "is_heading": is_heading,
                "is_bullet": is_bullet,
            })
        if paragraphs:
            slides.append(paragraphs)
    return _merge_orphan_headings(slides), []


def _parse_docx(path: Path) -> tuple[list[list[dict]], list[str]]:
    """
    Parse a DOCX file using blank paragraphs as slide boundaries.
    Returns (slides, warnings).  Images and tables are skipped with warnings.
    """
    from docx import Document  # imported here so PDF-only paths don't need it
    doc = Document(str(path))

    warnings = []
    slides: list[list[dict]] = []
    current: list[dict] = []

    # Warn once about tables
    if doc.tables:
        warnings.append(
            f"{len(doc.tables)} table(s) skipped — tables are not supported in DOCX conversion."
        )

    consecutive_blank = 0
    _BULLET_CHARS = ('•', '-', '*', '·', '◦', '▪', '▸', '►')

    for para in doc.paragraphs:
        style_name = (para.style.name or "").lower() if para.style else ""

        # --- Multi-line paragraph (Shift+Enter / <w:br> line breaks) ----------
        # Some DOCX files pack an entire slide's content into one <w:p> using
        # soft line breaks, with literal • characters as bullet markers.
        # Detect this pattern and split into a self-contained slide.
        raw_text = para.text.replace("\xa0", " ")
        if "\n" in raw_text:
            lines = [l.strip() for l in raw_text.split("\n") if l.strip()]
            if lines:
                # Flush any accumulated content as its own slide first
                if current:
                    slides.append(current)
                    current = []
                consecutive_blank = 0

                slide_paras = []
                found_heading = False
                for line in lines:
                    if line.startswith(_BULLET_CHARS):
                        stripped = line.lstrip(''.join(_BULLET_CHARS)).strip()
                        if stripped:
                            slide_paras.append({
                                "text": stripped,
                                "runs": [{"text": stripped, "bold": False, "italic": False,
                                          "underline": False, "color": None, "highlight": None}],
                                "is_heading": False,
                                "is_bullet": True,
                            })
                    else:
                        slide_paras.append({
                            "text": line,
                            "runs": [{"text": line, "bold": False, "italic": False,
                                      "underline": False, "color": None, "highlight": None}],
                            "is_heading": not found_heading,
                            "is_bullet": False,
                        })
                        found_heading = True

                if slide_paras:
                    slides.append(slide_paras)
            continue
        # --- End multi-line paragraph handling --------------------------------

        is_heading = "heading" in style_name
        is_bullet  = "list" in style_name or "bullet" in style_name

        # Fallback: check OOXML numPr element — bullets from Google Docs,
        # Keynote, and other exporters often carry no named list style but
        # always have a <w:numPr> block in the paragraph XML.
        if not is_bullet:
            _WML_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            # ⚠ DO NOT change `{{{...}}}` to `{{...}}` — triple braces produce a
            #   LITERAL `{` after f-string substitution (Clark notation); double
            #   braces inline the namespace value and lxml rejects it. See docs/DOCX_TXT.md.
            is_bullet = para._p.find(f".//{{{_WML_NS}}}numPr") is not None

        # Blank paragraph — count them; two or more in a row = slide boundary
        if not para.text.strip():
            consecutive_blank += 1
            if consecutive_blank >= 2 and current:
                slides.append(current)
                current = []
            continue

        consecutive_blank = 0

        # Build run list, preserving all character formatting
        _WML = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        runs = []
        has_image = False
        for run in para.runs:
            if run.element.find(f".//{{{_WML}}}drawing") is not None:
                has_image = True

            if not run.text:
                continue

            # Explicit RGB color (ignore theme/auto colors — they're not portable)
            color = None
            try:
                from docx.enum.dml import MSO_COLOR_TYPE
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    rgb = run.font.color.rgb   # python-docx RGBColor is a (r,g,b) tuple
                    color = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            except Exception:
                pass

            # Highlight color (background highlight like yellow marker)
            highlight = None
            try:
                hl = run.font.highlight_color
                if hl is not None:
                    # python-docx returns WD_COLOR_INDEX enum; map common ones to RGB
                    _HL_MAP = {
                        1:  (255, 255,   0),   # yellow
                        2:  (0,   255, 255),   # cyan
                        3:  (255,   0, 255),   # magenta
                        4:  (0,   255,   0),   # bright green
                        5:  (0,     0, 255),   # blue
                        6:  (255,   0,   0),   # red
                        7:  (0,     0, 128),   # dark blue
                        8:  (0,   128, 128),   # teal
                        9:  (0,   128,   0),   # green
                        10: (128,   0, 128),   # dark magenta
                        11: (128,   0,   0),   # dark red
                        12: (128, 128,   0),   # dark yellow
                        14: (192, 192, 192),   # light gray
                        15: (128, 128, 128),   # dark gray
                    }
                    highlight = _HL_MAP.get(int(hl))
            except Exception:
                pass

            runs.append({
                "text":      run.text,
                "bold":      bool(run.bold),
                "italic":    bool(run.italic),
                "underline": bool(run.underline),
                "color":     color,      # (r,g,b) or None
                "highlight": highlight,  # (r,g,b) or None
            })

        if has_image:
            warnings.append(
                f"Slide {len(slides) + 1}: inline image skipped."
            )

        if not runs:
            runs = [{"text": para.text, "bold": False, "italic": False,
                     "underline": False, "color": None, "highlight": None}]

        current.append({
            "text": para.text,
            "runs": runs,
            "is_heading": is_heading,
            "is_bullet": is_bullet,
        })

    if current:
        slides.append(current)

    return _merge_orphan_headings(slides), warnings


# --- Font-fitting ---

def _estimate_fits(paragraphs: list[dict], font_size: float) -> bool:
    """Estimate whether paragraph list fits on a slide at font_size (pts)."""
    cpl = _TXT_W_PT / (font_size * _CHAR_W_RATIO)   # chars per line
    lines_avail = _TXT_H_PT / (font_size * _LINE_H_RATIO)

    total = 0.0
    for i, para in enumerate(paragraphs):
        text = ("• " if para.get("is_bullet") else "") + para.get("text", "")
        if not text.strip():
            total += 0.5
            continue
        eff = font_size * 1.2 if para.get("is_heading") else font_size
        eff_cpl = _TXT_W_PT / (eff * _CHAR_W_RATIO)
        total += max(1.0, math.ceil(len(text) / eff_cpl))
        if i < len(paragraphs) - 1:
            total += 0.4   # inter-paragraph gap
    return total <= lines_avail


def _fit_font_size(paragraphs: list[dict]) -> tuple[int, bool]:
    """Return (font_size, overflowed). Binary search _MIN_FONT_PT.._MAX_FONT_PT."""
    if _estimate_fits(paragraphs, _MAX_FONT_PT):
        return _MAX_FONT_PT, False
    lo, hi = _MIN_FONT_PT, _MAX_FONT_PT
    while lo < hi:
        mid = (lo + hi + 1) // 2
        if _estimate_fits(paragraphs, mid):
            lo = mid
        else:
            hi = mid - 1
    return lo, not _estimate_fits(paragraphs, lo)


# --- Slide builder ---

def _add_text_slide(
    prs: Presentation,
    paragraphs: list[dict],
    base_size: int,
    dark_mode: bool = False,
    text_align: str = "center",
) -> None:
    """Add one text slide to the presentation."""
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)

    if dark_mode:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

    side_margin = Inches(_TXT_MARGIN_IN)
    top_margin  = Inches(_TXT_TOP_MARGIN_IN)
    txBox  = slide.shapes.add_textbox(
        left   = side_margin,
        top    = top_margin,
        width  = Emu(SLIDE_WIDTH_EMU)  - side_margin * 2,
        height = Emu(SLIDE_HEIGHT_EMU) - top_margin  * 2,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    align = PP_ALIGN.LEFT if text_align == "left" else PP_ALIGN.CENTER

    for i, para_data in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align

        is_heading = para_data.get("is_heading", False)
        is_bullet  = para_data.get("is_bullet",  False)
        eff_size   = int(base_size * 1.2) if is_heading else base_size

        runs = para_data.get("runs") or [{"text": para_data.get("text", ""), "bold": False, "italic": False}]

        if is_bullet:
            br = p.add_run()
            br.text       = "• "
            br.font.size  = Pt(eff_size)
            if dark_mode:
                br.font.color.rgb = RGBColor(255, 255, 255)

        for rd in runs:
            r = p.add_run()
            r.text        = rd.get("text", "")
            r.font.size   = Pt(eff_size)
            r.font.bold   = True if is_heading else (rd.get("bold") or False)
            r.font.italic = rd.get("italic")    or False
            if rd.get("underline"):
                r.font.underline = True
            if rd.get("color"):
                # Explicit color from source doc always wins, even in dark mode
                r.font.color.rgb = RGBColor(*rd["color"])
            elif dark_mode:
                r.font.color.rgb = RGBColor(255, 255, 255)


# --- Overflow analysis ---

def docx_analyze(path: Path) -> dict:
    """Parse a DOCX/TXT and report which slides overflow at minimum font size."""
    ext = path.suffix.lower()
    try:
        if ext == ".txt":
            slides, _ = _parse_txt(path)
        elif ext == ".docx":
            slides, _ = _parse_docx(path)
        else:
            return {"ok": False, "error": f"Unsupported file type: {ext}"}
    except Exception as e:
        return {"ok": False, "error": str(e)}

    overflow_slides = []
    for i, paragraphs in enumerate(slides):
        _, overflowed = _fit_font_size(paragraphs)
        if overflowed:
            heading = next(
                (p["text"].strip() for p in paragraphs if p.get("is_heading") and p.get("text", "").strip()),
                None,
            )
            lines = [p["text"] for p in paragraphs if p.get("text", "").strip()]
            overflow_slides.append({
                "slide_index": i,
                "heading": heading,
                "lines": lines,
            })

    return {"ok": True, "total_slides": len(slides), "overflow_slides": overflow_slides}


def _find_spill_break(paragraphs: list[dict]) -> int:
    """Return index to split after: largest k where paragraphs[:k] fits at min font."""
    lo, hi = 1, len(paragraphs) - 1
    while lo < hi:
        mid = (lo + hi + 1) // 2
        _, overflowed = _fit_font_size(paragraphs[:mid])
        if not overflowed:
            lo = mid
        else:
            hi = mid - 1
    return max(1, lo)


def _find_midpoint_break(paragraphs: list[dict]) -> int:
    """Split near midpoint, never ending a slide mid-sentence."""
    n = len(paragraphs)
    if n <= 1:
        return 1
    SENTENCE_ENDS = {'.', '!', '?', '…'}
    mid = n // 2
    for delta in range(n):
        for k in [mid - delta, mid + delta]:
            if 0 < k < n:
                text = paragraphs[k - 1].get("text", "").rstrip()
                if text and text[-1] in SENTENCE_ENDS:
                    return k
    return mid


def _apply_split_choices(slides: list, choices: list) -> list:
    """Apply user split choices to the parsed slide list. Returns expanded list."""
    import json as _json  # already imported at module level but safe to re-import

    # Process in reverse order so index offsets don't cascade
    result = list(slides)
    for choice in sorted(choices, key=lambda c: c.get("slide_index", 0), reverse=True):
        idx = choice.get("slide_index", 0)
        if idx >= len(result):
            continue
        paragraphs = result[idx]
        mode = choice.get("mode", "midpoint")

        if mode == "spill":
            break_at = _find_spill_break(paragraphs)
        elif mode == "midpoint":
            break_at = _find_midpoint_break(paragraphs)
        elif mode == "manual":
            break_at = int(choice.get("break_line", _find_midpoint_break(paragraphs)))
        else:
            continue

        break_at = max(1, min(break_at, len(paragraphs) - 1))
        part1 = paragraphs[:break_at]
        part2 = list(paragraphs[break_at:])

        if not part2:
            continue

        # Prepend cont. heading if first paragraph is a heading
        heading_text = None
        if paragraphs and paragraphs[0].get("is_heading"):
            heading_text = paragraphs[0].get("text", "").strip()
        if heading_text:
            cont_text = f"{heading_text} cont."
            part2.insert(0, {
                "text": cont_text,
                "runs": [{"text": cont_text, "bold": True, "italic": False,
                          "underline": False, "color": None, "highlight": None}],
                "is_heading": True,
                "is_bullet": False,
            })

        result[idx] = part1
        result.insert(idx + 1, part2)

    return result


# --- docx_info (for IPC query before conversion) ---

def docx_info(path: Path) -> dict:
    """
    Count slides and words in a .txt or .docx file.
    Emits {"type": "docx_info", ...} in IPC mode.
    """
    ext = path.suffix.lower()
    try:
        if ext == ".txt":
            slides, _ = _parse_txt(path)
        elif ext == ".docx":
            slides, _ = _parse_docx(path)
        else:
            raise ValueError(f"Unsupported extension: {ext}")
        word_count = sum(
            len(p["text"].split()) for slide in slides for p in slide
        )
        result = {"ok": True, "slideCount": len(slides), "wordCount": word_count}
    except Exception as e:
        result = {"ok": False, "message": str(e)}

    if _ipc_mode:
        print(json.dumps({"type": "docx_info", **result}), flush=True)
    return result


# --- Main converter ---

def convert_text_doc(
    file_path: Path,
    output_dir: Path,
    suffix: str = "",
    file_index: int = 1,
    total_files: int = 1,
    slide_theme: str = "light",
    text_align: str = "left",
    split_choices: list | None = None,
) -> dict:
    """Convert a .txt or .docx file to a PPTX using blank-line slide boundaries."""
    ext  = file_path.suffix.lower()
    stem = file_path.stem + suffix
    out_path = output_dir / f"{stem}.pptx"

    _emit({
        "type": "start",
        "file": str(file_path),
        "file_index": file_index,
        "total_files": total_files,
    })

    try:
        if ext == ".txt":
            slides_data, parse_warnings = _parse_txt(file_path)
        else:
            slides_data, parse_warnings = _parse_docx(file_path)
    except Exception as e:
        msg = f"Failed to parse {file_path.name}: {e}"
        _emit({"type": "error", "file": str(file_path), "message": msg})
        return {"ok": False, "message": msg}

    for w in parse_warnings:
        _emit({"type": "warn", "file": str(file_path), "message": w})

    if not slides_data:
        msg = "No slide content found — is the file empty?"
        _emit({"type": "error", "file": str(file_path), "message": msg})
        return {"ok": False, "message": msg}

    # Apply user split choices (overflow resolution from UI)
    if split_choices:
        slides_data = _apply_split_choices(slides_data, split_choices)

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    total_slides = len(slides_data)

    for i, paragraphs in enumerate(slides_data, start=1):
        _emit({
            "type": "progress",
            "file": str(file_path),
            "page": i,
            "total_pages": total_slides,
            "message": f"Building slide {i} of {total_slides}",
        })

        font_size, overflowed = _fit_font_size(paragraphs)

        if overflowed:
            _emit({
                "type": "warn",
                "file": str(file_path),
                "message": (
                    f"Slide {i}: text overflows at {_MIN_FONT_PT}pt — "
                    "some text may be cut off. Consider splitting this block."
                ),
            })

        try:
            _add_text_slide(prs, paragraphs, font_size, dark_mode=(slide_theme == "dark"), text_align=text_align)
        except Exception as e:
            msg = f"Failed building slide {i}: {e}"
            _emit({"type": "error", "file": str(file_path), "message": msg})
            return {"ok": False, "message": msg}

    try:
        prs.save(str(out_path))
    except (PermissionError, OSError) as e:
        msg = f"Cannot save {out_path.name}: {e}"
        _emit({"type": "error", "file": str(file_path), "message": msg})
        return {"ok": False, "message": msg}

    _emit({
        "type": "done",
        "file": str(file_path),
        "output": str(out_path),
        "slides": total_slides,
    })
    return {"ok": True, "output": str(out_path), "slides": total_slides}


# ---------------------------------------------------------------------------
# PPTX → notes feature (Mode A: Teleprompter, Mode B: Split Deck, Mode C: Companion)
# ---------------------------------------------------------------------------

_PPTX_SHAPE_TAGS = {qn('p:sp'), qn('p:pic'), qn('p:graphicFrame'),
                    qn('p:grpSp'), qn('p:cxnSp'), qn('p:contentPart')}
_R_NS_BRACE = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_PRES_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_PPTX_MARKER_RE = re.compile(r"^-{3,}\s*$")


# --- Generic helpers (used by all three modes) ---

def _pptx_get_notes(slide) -> str:
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text or ""
    return ""


def _pptx_set_notes(slide, text: str) -> None:
    """Write text to a slide's notes. Lazily creates the notes_slide if needed."""
    slide.notes_slide.notes_text_frame.text = text


def _remap_rids(root, rid_map: dict) -> None:
    """Rewrite every relationship-namespace attribute in `root` and descendants."""
    for el in root.iter():
        for attr_name in list(el.attrib.keys()):
            if attr_name.startswith(_R_NS_BRACE):
                old = el.attrib[attr_name]
                if old in rid_map:
                    el.attrib[attr_name] = rid_map[old]


def _pptx_clone_slide(prs, src_slide, keep_animations: bool = True):
    """Clone a source slide into prs.slides (appends to end). Returns new slide.

    Rels-rebind + rId-remap recipe — preserves images, hyperlinks, embedded
    objects, transitions. Animations (timing) preserved when keep_animations=True;
    stripped on subsequent clones in Mode B/C so click-builds only fire on the
    first instance of a duplicated slide.
    """
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    # Strip placeholder shapes the layout injected
    new_sp_tree = new_slide.shapes._spTree
    for sp in list(new_sp_tree):
        if sp.tag in _PPTX_SHAPE_TAGS:
            new_sp_tree.remove(sp)

    # Build src_rId → new_rId map (skip notesSlide — new slide gets fresh notes).
    # External rels (URL hyperlinks) use get_or_add_ext_rel; owned target parts
    # (charts, SmartArt diagram bundle, OLE embeds) get a recursive deep copy
    # so the clone owns its own resources and PowerPoint stays happy.
    rid_map: dict = {}
    for rel in src_slide.part.rels.values():
        # ⚠ DO NOT remove — clones MUST get a fresh notes pane; carrying the
        #   source's notesSlide rel causes shared-notes corruption between
        #   duplicates. See docs/PPTX_PIPELINE.md.
        if "notesSlide" in rel.reltype:
            continue
        # ⚠ DO NOT skip is_external — rel.target_part is undefined for external
        #   rels (URL hyperlinks) and accessing it raises. See docs/PPTX_PIPELINE.md.
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        elif rel.reltype in _PPTX_OWNED_RELTYPES:
            owned_copy = _pptx_deep_copy_part(prs.part.package, rel.target_part)
            new_rid = new_slide.part.rels.get_or_add(rel.reltype, owned_copy)
        else:
            new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[rel.rId] = new_rid

    # Deep-copy shapes, rewriting rIds in the copied subtree
    for sp in src_slide.shapes._spTree:
        if sp.tag in _PPTX_SHAPE_TAGS:
            new_sp = copy.deepcopy(sp)
            _remap_rids(new_sp, rid_map)
            new_sp_tree.append(new_sp)

    src_sld = src_slide._element
    new_sld = new_slide._element
    src_cSld = src_sld.find(qn('p:cSld'))
    new_cSld = new_sld.find(qn('p:cSld'))

    # Carry over slide-level <p:bg> (background image / fill at the slide level —
    # overrides the layout's background). Title slides often store their visual
    # entirely here, with zero shapes. Without this copy, clones of such slides
    # render blank.
    src_bg = src_cSld.find(qn('p:bg'))
    if src_bg is not None:
        for bg in new_cSld.findall(qn('p:bg')):
            new_cSld.remove(bg)
        new_bg = copy.deepcopy(src_bg)
        _remap_rids(new_bg, rid_map)
        new_cSld.insert(0, new_bg)  # <p:bg> comes before <p:spTree> per the schema

    # Carry over slide-level <p:clrMapOvr> (color map override; lives on <p:sld>,
    # sibling of <p:cSld>). Affects color theme on the slide.
    src_clr = src_sld.find(qn('p:clrMapOvr'))
    if src_clr is not None:
        for clr in new_sld.findall(qn('p:clrMapOvr')):
            new_sld.remove(clr)
        new_clr = copy.deepcopy(src_clr)
        _remap_rids(new_clr, rid_map)
        # Insert right after <p:cSld>
        new_sld.insert(list(new_sld).index(new_cSld) + 1, new_clr)

    # Carry over <p:transition> and (conditionally) <p:timing>
    src_transition = src_sld.find(qn('p:transition'))
    if src_transition is not None:
        for t in new_sld.findall(qn('p:transition')):
            new_sld.remove(t)
        new_transition = copy.deepcopy(src_transition)
        _remap_rids(new_transition, rid_map)
        new_sld.append(new_transition)

    if keep_animations:
        src_timing = src_sld.find(qn('p:timing'))
        if src_timing is not None:
            for t in new_sld.findall(qn('p:timing')):
                new_sld.remove(t)
            new_timing = copy.deepcopy(src_timing)
            _remap_rids(new_timing, rid_map)
            new_sld.append(new_timing)

    return new_slide


def _pptx_move_slide_to(prs, from_idx: int, to_idx: int) -> None:
    """Move a slide within prs.slides by manipulating <p:sldIdLst> directly."""
    sld_id_lst = prs.slides._sldIdLst
    children = list(sld_id_lst)
    child = children[from_idx]
    sld_id_lst.remove(child)
    sld_id_lst.insert(to_idx, child)


# Reltypes whose target parts can't be safely *shared* across cloned slides —
# PowerPoint's data model expects each chart / SmartArt diagram / OLE-embed to
# be owned by a single slide. The clone must get its own deep copy of these
# parts (and recursively any sub-parts they own, e.g. chart→embedded Excel,
# diagramData→diagramLayout).
_PPTX_OWNED_RELTYPES = frozenset(
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/" + t
    for t in (
        "chart",
        "diagramData", "diagramLayout", "diagramColors",
        "diagramQuickStyle", "diagramDrawing",
        "oleObject", "package",
    )
)


def _pptx_next_partname_like(package, src_partname, reserved=None) -> object:
    """Given an existing partname (PackURI) like /ppt/diagrams/data1.xml,
    return a free PackURI in the same family (/ppt/diagrams/data2.xml, etc.).

    `reserved` is a set of PackURIs already chosen during in-flight recursive
    deep-copies — these aren't yet reachable from the package's rel graph (so
    iter_parts() doesn't see them), but we still need to avoid colliding."""
    from pptx.opc.packuri import PackURI
    s = str(src_partname)
    m = re.match(r'^(.+?)(\d+)(\.[^.]+)$', s)
    if m:
        template = f"{m.group(1)}%d{m.group(3)}"
    else:
        base, _, ext = s.rpartition('.')
        template = f"{base}_copy%d.{ext}"

    used = {p.partname for p in package.iter_parts()}
    if reserved:
        used |= reserved
    n = 1
    while PackURI(template % n) in used:
        n += 1
    return PackURI(template % n)


def _pptx_deep_copy_part(package, src_part, _memo=None, _reserved=None):
    """Recursively duplicate a Part (and any of its owned sub-parts).
    Returns the new Part — caller is responsible for wiring it into the
    package's rel graph (which is what makes it persistable on save).

    Shared (non-owned) rel targets — images, hyperlinks, theme, slide layouts —
    stay shared with the source. Only reltypes in _PPTX_OWNED_RELTYPES trigger
    a recursive deep-copy. After the new part's rels are built, any rIds that
    were remapped get rewritten inside the new part's XML element so internal
    references resolve correctly.
    """
    if _memo is None:
        _memo = {}
    # ⚠ DO NOT remove _reserved — parts created inside the same recursion
    #   aren't reachable via iter_parts() yet, so the next-free-partname scan
    #   misses them and collides. See docs/PPTX_PIPELINE.md.
    if _reserved is None:
        _reserved = set()
    if src_part in _memo:
        return _memo[src_part]

    new_partname = _pptx_next_partname_like(package, src_part.partname, _reserved)
    _reserved.add(new_partname)
    new_part = type(src_part).load(
        new_partname, src_part.content_type, package, src_part.blob,
    )
    _memo[src_part] = new_part

    # Recreate the part's relationships; deep-copy any owned targets.
    rid_map: dict = {}
    for rel in src_part.rels.values():
        if rel.is_external:
            new_rid = new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        elif rel.reltype in _PPTX_OWNED_RELTYPES:
            sub_target = _pptx_deep_copy_part(package, rel.target_part, _memo, _reserved)
            new_rid = new_part.rels.get_or_add(rel.reltype, sub_target)
        else:
            new_rid = new_part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[rel.rId] = new_rid

    # If the part has an XML body, rewrite its internal rId references so the
    # new part's relationships are the ones in scope (diagram data XML
    # references its sibling layout/colors/etc. via rId from inside).
    el = getattr(new_part, "_element", None)
    if el is not None and rid_map:
        _remap_rids(el, rid_map)

    return new_part


def _pptx_count_click_builds(slide) -> int:
    """Count click-triggered build animations on a slide."""
    timing = slide._element.find(f".//{{{_PRES_NS}}}timing")
    if timing is None:
        return 0
    return sum(
        1 for ctn in timing.iter(f"{{{_PRES_NS}}}cTn")
        if ctn.get("nodeType") in ("clickEffect", "clickPar")
    )


# Threshold is now expressed in *effective lines* (vertical space in the
# notes pane), not characters. Lines map to actual display: every '\n' is at
# least one line, long lines wrap, blanks count as gaps. A line of 800 chars
# of prose costs the same vertical space as a line of 10 chars.
_PPTX_NOTES_CHARS_PER_LINE = 65   # typical width of one visual line in Presenter View notes
_PPTX_SOFT_BUFFER = 1.2           # ⚠ DO NOT tune outside 1.2–1.4 without re-checking the 13-line slide case. See docs/PPTX_PIPELINE.md.


def _effective_lines(text: str) -> int:
    """Vertical lines a piece of notes text would occupy when rendered.
    Counts each '\\n' as ≥1 line; wraps long lines; blank lines = 1-line gap."""
    if not (text and text.strip()):
        return 0
    cpl = _PPTX_NOTES_CHARS_PER_LINE
    total = 0
    for raw_line in text.split("\n"):
        line = raw_line.rstrip()
        if not line:
            total += 1
        else:
            total += max(1, (len(line) + cpl - 1) // cpl)
    return total


def _greedy_pack(parts: list[str], threshold: int, sep: str) -> list[str]:
    """Pack parts into chunks of ~threshold *effective lines*, never splitting a
    part. After the greedy pass, merge any chunk smaller than half-threshold
    into a neighbor — speaker names ("Joel", "Matt") and short closing lines
    that would otherwise become near-empty slides get absorbed.
    """
    chunks: list[str] = []
    current = ""
    for p in parts:
        candidate = (current + sep + p) if current else p
        if _effective_lines(candidate) <= threshold or not current:
            current = candidate
        else:
            chunks.append(current)
            current = p
    if current:
        chunks.append(current)

    if len(chunks) <= 1:
        return chunks

    min_lines = max(threshold // 2, 1)
    merged: list[str] = [chunks[0]]
    for c in chunks[1:]:
        if _effective_lines(c) < min_lines or _effective_lines(merged[-1]) < min_lines:
            merged[-1] = merged[-1] + sep + c
        else:
            merged.append(c)
    return merged


def _split_by_threshold(text: str, threshold: int) -> list[str]:
    """Sub-split a segment that exceeds threshold lines.

    Prefer paragraph-level packing only when every paragraph itself fits in
    threshold lines; otherwise fall through to sentence-level on the whole
    text (avoids splitting a tall paragraph badly across chunk boundaries).
    """
    paragraphs = [p.strip() for p in re.split(r"\n[ \t]*\n", text) if p.strip()]
    if len(paragraphs) > 1 and max(_effective_lines(p) for p in paragraphs) <= threshold:
        return _greedy_pack(paragraphs, threshold, "\n\n")

    sentences = re.split(r"(?<=[.!?…])\s+", text)
    if len(sentences) > 1:
        return _greedy_pack(sentences, threshold, " ")

    return _greedy_pack(text.split(), threshold, " ")


def _chunk_notes(text: str, threshold: int) -> list[str]:
    """Split notes into chunks. Override markers (--- on own line) always apply;
    threshold-based sub-splitting runs within each segment when effective lines
    exceed the soft buffer (1.2× threshold)."""
    text = (text or "").strip()
    if not text:
        return []

    # First pass: split on override markers (always active)
    segments: list[list[str]] = [[]]
    for line in text.split("\n"):
        if _PPTX_MARKER_RE.match(line.strip()):
            if segments[-1]:
                segments.append([])
            continue
        segments[-1].append(line)
    raw_segments = [
        "\n".join(seg).strip()
        for seg in segments
        if any(l.strip() for l in seg)
    ]

    # Second pass: threshold-split each segment if it would overflow vertically
    soft_limit = threshold * _PPTX_SOFT_BUFFER
    chunks: list[str] = []
    for seg in raw_segments:
        if _effective_lines(seg) <= soft_limit:
            chunks.append(seg)
        else:
            chunks.extend(_split_by_threshold(seg, threshold))
    return chunks


# --- Mode B helpers: notes decoration + per-slide chunk → clones ---

def _decorate_chunk(text: str, src_idx: int, chunk_idx: int, total: int,
                    include_slide_numbers: bool, append_continue: bool) -> str:
    """Apply optional slide-number prefix and [click to continue] suffix."""
    out = text
    if append_continue and chunk_idx < total - 1:
        out = out + "\n\n[click to continue]"
    if include_slide_numbers:
        if total == 1:
            prefix = f"[Slide {src_idx + 1}]"
        else:
            prefix = f"[Slide {src_idx + 1} — part {chunk_idx + 1}/{total}]"
        out = prefix + "\n" + out
    return out


def _atomic_save_pptx(prs, out_path: Path) -> None:
    """Save a presentation atomically — writes to <out>.tmp first then renames,
    so the final path never exists on disk in a partially-written state.
    PowerPoint trying to open a half-saved file is what gave us the
    WC23 'mid-conversion' bug; this prevents that class of failure."""
    tmp = out_path.with_name(out_path.name + ".tmp")
    prs.save(str(tmp))
    os.replace(str(tmp), str(out_path))


def _atomic_write_text(out_path: Path, content: str) -> None:
    tmp = out_path.with_name(out_path.name + ".tmp")
    tmp.write_text(content, encoding="utf-8")
    os.replace(str(tmp), str(out_path))


def _pptx_resolve_chunks(slide, threshold: int) -> tuple[list[str], bool]:
    """Compute chunks for a slide. Returns (chunks, was_merged); was_merged is
    always False now that _pptx_clone_slide can deep-copy owned target parts
    (charts, SmartArt, OLE embeds) — every slide can be safely duplicated.
    The tuple shape is preserved for caller compatibility."""
    notes = _pptx_get_notes(slide).strip()
    raw_chunks = _chunk_notes(notes, threshold) if notes else []
    return raw_chunks, False


def _pptx_split_deck_into(
    src_path: Path,
    out_path: Path,
    threshold: int,
    include_slide_numbers: bool,
) -> tuple[int, list[dict]]:
    """Build a split-deck variant of src_path at out_path.

    Returns (final_slide_count, edits) — where `edits` is a list of records
    describing per-source-slide changes (duplicated, merged_uncloneable) for
    the audit summary. Slides that passed through unchanged are not included.
    """
    prs = Presentation(str(src_path))
    original_slides = list(prs.slides)
    target_idx = 0
    edits: list[dict] = []

    for src_idx, src_slide in enumerate(original_slides):
        chunks, was_merged = _pptx_resolve_chunks(src_slide, threshold)
        src_slide_no = src_idx + 1

        if was_merged:
            _emit({
                "type": "warn",
                "file": str(src_path),
                "message": (
                    f"Slide {src_slide_no}: contains a chart, SmartArt diagram, "
                    "or embedded object — notes can't be split across duplicate slides. "
                    "All chunks merged into the one slide's notes pane."
                ),
            })
            edits.append({
                "src_slide": src_slide_no,
                "kind": "merged_uncloneable",
                "result_slides": [target_idx + 1],
            })

        if len(chunks) <= 1:
            if chunks:
                _pptx_set_notes(src_slide, _decorate_chunk(
                    chunks[0], src_idx, 0, 1, include_slide_numbers, False,
                ))
            target_idx += 1
            continue

        # First chunk → source slide (preserves its animations)
        _pptx_set_notes(src_slide, _decorate_chunk(
            chunks[0], src_idx, 0, len(chunks), include_slide_numbers, True,
        ))
        result_slides_for_this = [target_idx + 1]
        target_idx += 1

        # Remaining chunks → clones (no animations) inserted right after source
        for k in range(1, len(chunks)):
            clone = _pptx_clone_slide(prs, src_slide, keep_animations=False)
            _pptx_set_notes(clone, _decorate_chunk(
                chunks[k], src_idx, k, len(chunks), include_slide_numbers, True,
            ))
            end_idx = len(prs.slides) - 1
            if end_idx != target_idx:
                _pptx_move_slide_to(prs, end_idx, target_idx)
            result_slides_for_this.append(target_idx + 1)
            target_idx += 1

        edits.append({
            "src_slide": src_slide_no,
            "kind": "duplicated",
            "chunk_count": len(chunks),
            "result_slides": result_slides_for_this,
        })

    _atomic_save_pptx(prs, out_path)
    return len(prs.slides), edits


def _pptx_split_deck(
    src_path: Path,
    output_dir: Path,
    threshold: int,
    include_slide_numbers: bool,
    suffix: str,
) -> dict:
    out_path = output_dir / f"{src_path.stem}_split{suffix}.pptx"
    slides, edits = _pptx_split_deck_into(src_path, out_path, threshold, include_slide_numbers)
    return {"ok": True, "output": str(out_path), "slides": slides, "edits": edits}


# --- Mode A: Teleprompter (extract notes to .txt) ---

def _pptx_teleprompter(
    src_path: Path,
    output_dir: Path,
    threshold: int,
    include_headers: bool,
    suffix: str,
) -> dict:
    prs = Presentation(str(src_path))
    out_path = output_dir / f"{src_path.stem}_notes{suffix}.txt"

    lines_out: list[str] = []
    edits: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        notes = _pptx_get_notes(slide).strip()
        # Empty notes get the placeholder so the speaker can see the slide is intentionally bare
        if not notes:
            chunks = ["[no notes]"]
            edits.append({"src_slide": i, "kind": "empty"})
        else:
            chunks = _chunk_notes(notes, threshold) or [notes]
            if len(chunks) > 1:
                edits.append({"src_slide": i, "kind": "chunked", "chunk_count": len(chunks)})

        for j, chunk in enumerate(chunks):
            if include_headers:
                if len(chunks) == 1:
                    lines_out.append(f"[Slide {i}]")
                else:
                    lines_out.append(f"[Slide {i} — part {j + 1}/{len(chunks)}]")
            else:
                # When headers are off, a blank line still separates slide-blocks
                if lines_out and lines_out[-1] != "":
                    lines_out.append("")
            lines_out.append(chunk)
            lines_out.append("")

    _atomic_write_text(out_path, "\n".join(lines_out).rstrip() + "\n")
    return {"ok": True, "output": str(out_path), "slides": len(prs.slides), "edits": edits}


# --- Mode C: Companion (paired audience + companion decks, click-synced) ---

def _build_companion_paragraphs(chunk_text: str, slide_number: int | None = None) -> list[dict]:
    """Build the paragraph list for one companion slide. Reuses the existing
    text-slide builder by emitting paragraph dicts in the same shape as
    _parse_txt() / _parse_docx() produce."""
    paragraphs: list[dict] = []
    if slide_number is not None:
        head = f"Slide {slide_number}"
        paragraphs.append({
            "text": head,
            "runs": [{"text": head, "bold": True, "italic": False,
                      "underline": False, "color": None, "highlight": None}],
            "is_heading": True,
            "is_bullet": False,
        })

    chunk_text = (chunk_text or "").strip() or "[no notes]"
    # Split per line — same convention as _parse_txt — so _fit_font_size counts
    # the actual rendered lines and the text-box auto-fit math stays accurate.
    raw_lines = [line.strip() for line in chunk_text.split("\n") if line.strip()]
    if not raw_lines:
        raw_lines = [chunk_text]

    for line in raw_lines:
        paragraphs.append({
            "text": line,
            "runs": [{"text": line, "bold": False, "italic": False,
                      "underline": False, "color": None, "highlight": None}],
            "is_heading": False,
            "is_bullet": False,
        })
    return paragraphs


def _pptx_companion(
    src_path: Path,
    output_dir: Path,
    threshold: int,
    slide_theme: str,
    include_slide_numbers: bool,
    suffix: str,
) -> dict:
    """Produce paired audience + companion decks. Click-synced: per source slide
    N with B builds and K chunks, audience side has K duplicates (first keeps
    builds, rest stripped), companion side has B+K slides:
        chunks[0] × (B + 1)    # covers all build clicks + advance off slide 1
        chunks[1..K-1] × 1     # each covers one audience duplicate
    """
    audience_path = output_dir / f"{src_path.stem}_audience{suffix}.pptx"
    companion_path = output_dir / f"{src_path.stem}_companion{suffix}.pptx"

    # Audience deck — same engine as Mode B, just different filename
    audience_slides, audience_edits = _pptx_split_deck_into(
        src_path, audience_path, threshold, include_slide_numbers,
    )

    # Companion deck — text-slides via the existing _add_text_slide engine.
    # Use the canonical 13.33×7.5 dimensions that _add_text_slide / _fit_font_size
    # expect; the companion is a standalone notes deck for a confidence monitor
    # and doesn't need to inherit the audience deck's aspect ratio.
    src_prs = Presentation(str(src_path))
    companion_prs = Presentation()
    companion_prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    companion_prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    total_companion = 0
    for i, src_slide in enumerate(src_prs.slides, start=1):
        chunks, _was_merged = _pptx_resolve_chunks(src_slide, threshold)
        if not chunks:
            chunks = ["[no notes]"]
        builds = _pptx_count_click_builds(src_slide)

        # ⚠ DO NOT simplify this formula. chunks[0] × (builds + 1) covers
        #   every build-click PLUS the advance off slide 1; chunks[1..] each
        #   cover one audience duplicate. This is the only shape that holds
        #   the click-sync invariant. See docs/PPTX_PIPELINE.md.
        expanded = [chunks[0]] * (builds + 1) + chunks[1:]

        for chunk in expanded:
            paragraphs = _build_companion_paragraphs(
                chunk, slide_number=(i if include_slide_numbers else None),
            )
            font_size, _ = _fit_font_size(paragraphs)
            _add_text_slide(
                companion_prs, paragraphs, font_size,
                dark_mode=(slide_theme == "dark"), text_align="left",
            )
            total_companion += 1

    _atomic_save_pptx(companion_prs, companion_path)

    return {
        "ok": True,
        "outputs": [str(audience_path), str(companion_path)],
        "slides": audience_slides,
        "companion_slides": total_companion,
        "edits": audience_edits,
    }


# --- Dispatcher ---

def convert_pptx_notes(
    file_path: Path,
    output_dir: Path,
    mode: str = "teleprompter",
    threshold: int = 12,
    include_headers: bool = True,
    include_slide_numbers: bool = False,
    slide_theme: str = "dark",
    suffix: str = "",
    file_index: int = 1,
    total_files: int = 1,
) -> dict:
    """Convert a .pptx file based on the selected mode."""
    _emit({
        "type": "start",
        "file": str(file_path),
        "file_index": file_index,
        "total_files": total_files,
    })

    try:
        if mode == "teleprompter":
            result = _pptx_teleprompter(
                file_path, output_dir, threshold, include_headers, suffix,
            )
        elif mode == "split":
            result = _pptx_split_deck(
                file_path, output_dir, threshold, include_slide_numbers, suffix,
            )
        elif mode == "companion":
            result = _pptx_companion(
                file_path, output_dir, threshold, slide_theme,
                include_slide_numbers, suffix,
            )
        else:
            msg = f"Unknown pptx mode: {mode}"
            _emit({"type": "error", "file": str(file_path), "message": msg})
            return {"ok": False, "message": msg}
    except Exception as e:
        import traceback
        traceback.print_exc(file=sys.stderr)  # surfaces in app log via main.js stderr handler
        msg = f"PPTX conversion failed: {e}"
        _emit({"type": "error", "file": str(file_path), "message": msg})
        return {"ok": False, "message": msg}

    done_payload = {
        "type": "done",
        "file": str(file_path),
        "output": result.get("output") or (result.get("outputs") or [""])[0],
        "slides": result.get("slides", 0),
    }
    if "outputs" in result:
        done_payload["outputs"] = result["outputs"]
    if "companion_slides" in result:
        done_payload["companion_slides"] = result["companion_slides"]
    if "edits" in result:
        done_payload["edits"] = result["edits"]
        done_payload["pptx_mode"] = mode
    _emit(done_payload)
    return result


# --- pptx_info (for IPC query before conversion) ---

def pptx_info(path: Path) -> dict:
    """Return slide count, notes stats, build counts. Emits JSON in IPC mode."""
    try:
        prs = Presentation(str(path))
        slides = list(prs.slides)
        notes_lens = [len(_pptx_get_notes(s)) for s in slides]
        build_counts = [_pptx_count_click_builds(s) for s in slides]
        result = {
            "ok": True,
            "slideCount": len(slides),
            "totalNotesChars": sum(notes_lens),
            "slidesWithNotes": sum(1 for n in notes_lens if n > 0),
            "totalBuilds": sum(build_counts),
            "slidesWithBuilds": sum(1 for b in build_counts if b > 0),
        }
    except Exception as e:
        result = {"ok": False, "message": str(e)}

    if _ipc_mode:
        print(json.dumps({"type": "pptx_info", **result}), flush=True)
    return result


# ---------------------------------------------------------------------------
# Preflight check (called by Diagnostics tab in Electron)
# ---------------------------------------------------------------------------

def run_preflight(poppler_path: str | None = None) -> dict:
    """
    Run a full health check. Returns a dict of check → {ok, message}.
    Always emits a JSON line with type "preflight_result" when in IPC mode.
    """
    results = {}

    # 1. Poppler
    try:
        import subprocess
        cmd = [os.path.join(poppler_path, "pdftoppm") if poppler_path else "pdftoppm", "-v"]
        proc = subprocess.run(cmd, capture_output=True, timeout=5)
        ver = (proc.stdout or proc.stderr).decode(errors="replace").strip().split("\n")[0]
        results["poppler"] = {"ok": True, "message": ver}
    except Exception as e:
        results["poppler"] = {"ok": False, "message": str(e)}

    # 2. python-pptx
    try:
        import pptx
        results["python_pptx"] = {"ok": True, "message": f"python-pptx {pptx.__version__}"}
    except Exception as e:
        results["python_pptx"] = {"ok": False, "message": str(e)}

    # 3. Pillow / pdf2image
    try:
        import PIL
        results["pillow"] = {"ok": True, "message": f"Pillow {PIL.__version__}"}
    except Exception as e:
        results["pillow"] = {"ok": False, "message": str(e)}

    try:
        import pdf2image
        # pdf2image doesn't reliably expose __version__; confirm import + pdfinfo callable
        from pdf2image import pdfinfo_from_path  # noqa: F401
        ver = getattr(pdf2image, "__version__", "installed")
        results["pdf2image"] = {"ok": True, "message": f"pdf2image {ver}"}
    except Exception as e:
        results["pdf2image"] = {"ok": False, "message": str(e)}

    # 4. Python version
    import platform
    results["python"] = {
        "ok": True,
        "message": f"Python {sys.version.split()[0]} on {platform.system()} {platform.machine()}"
    }

    if _ipc_mode:
        print(json.dumps({"type": "preflight_result", "results": results}), flush=True)

    return results


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    global _ipc_mode

    parser = argparse.ArgumentParser(
        description="SlideFluid 3.0 — PDF to PPTX conversion engine",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("inputs", nargs="*", help="PDF files or folders")
    parser.add_argument("--output-dir", "-o", default=None,
                        help="Output folder (default: same as each input file)")
    parser.add_argument("--dpi", type=int, choices=[72, 144], default=72,
                        help="Raster DPI (default: 72)")
    parser.add_argument("--fill", choices=["black", "color_match", "smear"],
                        default="black", help="Pillarbox fill mode (default: black)")
    parser.add_argument("--text-align", choices=["left", "center"], default="left",
                        help="Text alignment for DOCX/TXT slides")
    parser.add_argument("--slide-theme", choices=["light", "dark"],
                        default="light", help="Slide background theme for text docs (default: light)")
    parser.add_argument("--overwrite", action="store_true",
                        help="Overwrite existing .pptx without asking")
    parser.add_argument("--skip-existing", action="store_true",
                        help="Skip files that already have a .pptx")
    parser.add_argument("--suffix", default="",
                        help="Append suffix to output filenames")
    parser.add_argument("--ipc", action="store_true",
                        help="Emit newline-delimited JSON (Electron IPC mode)")
    parser.add_argument("--poppler-path", default=None,
                        help="Path to Poppler bin directory (for bundled binaries)")
    parser.add_argument("--preflight", action="store_true",
                        help="Run preflight health check and exit")
    parser.add_argument("--docx-info", default=None, metavar="FILE",
                        help="Return slide/word count for a .txt or .docx file (IPC mode)")
    parser.add_argument("--analyze", default=None, metavar="FILE",
                        help="Analyze a DOCX/TXT for overflow slides and emit JSON")
    parser.add_argument("--split-choices", default=None, metavar="JSON",
                        help="JSON array of split choices [{slide_index,mode,break_line}]")
    parser.add_argument("--pptx-mode", choices=["teleprompter", "split", "companion"],
                        default="teleprompter",
                        help="PPTX input mode (default: teleprompter)")
    parser.add_argument("--pptx-threshold", type=int, default=12,
                        help="Max effective lines per notes chunk (default: 12)")
    parser.add_argument("--pptx-include-headers", choices=["yes", "no"], default="yes",
                        help="Include [Slide N] headers in teleprompter output")
    parser.add_argument("--pptx-slide-numbers", choices=["yes", "no"], default="no",
                        help="Include slide numbers on split/companion slides")
    parser.add_argument("--pptx-theme", choices=["light", "dark"], default="dark",
                        help="Companion deck theme (default: dark)")
    parser.add_argument("--pptx-info", default=None, metavar="FILE",
                        help="Return slide/notes/build stats for a .pptx (IPC mode)")

    args = parser.parse_args()
    _ipc_mode = args.ipc

    if args.docx_info:
        _ipc_mode = True  # always emit JSON for this command
        docx_info(Path(args.docx_info))
        sys.exit(0)

    if args.analyze:
        _ipc_mode = True
        result = docx_analyze(Path(args.analyze))
        print(json.dumps({"type": "docx_analyze", **result}))
        sys.exit(0)

    if args.pptx_info:
        _ipc_mode = True
        pptx_info(Path(args.pptx_info))
        sys.exit(0)

    if args.preflight:
        results = run_preflight(args.poppler_path)
        if not _ipc_mode:
            print("\nPreflight check results:")
            for check, r in results.items():
                status = "✓" if r["ok"] else "✗"
                print(f"  {status} {check:15s} {r['message']}")
        sys.exit(0 if all(r["ok"] for r in results.values()) else 1)

    if not args.inputs:
        parser.print_help()
        sys.exit(0)

    pdfs = collect_files(args.inputs)
    if not pdfs:
        print("No supported files found.", file=sys.stderr)
        sys.exit(1)

    if not _ipc_mode:
        print(f"SlideFluid 3.0 — {len(pdfs)} file(s) queued")
        print(f"  DPI: {args.dpi}  |  Fill: {args.fill}  |  Suffix: '{args.suffix}'")

    split_choices = json.loads(args.split_choices) if args.split_choices else None

    sys.exit(
        run_batch(
            pdf_paths=pdfs,
            output_dir=Path(args.output_dir) if args.output_dir else None,
            dpi=args.dpi,
            fill_mode=args.fill,
            overwrite=args.overwrite,
            skip_existing=args.skip_existing,
            suffix=args.suffix,
            poppler_path=args.poppler_path,
            slide_theme=args.slide_theme,
            text_align=args.text_align,
            split_choices=split_choices,
            pptx_opts={
                "mode": args.pptx_mode,
                "threshold": args.pptx_threshold,
                "include_headers": args.pptx_include_headers == "yes",
                "include_slide_numbers": args.pptx_slide_numbers == "yes",
                "slide_theme": args.pptx_theme,
            },
        )
    )


# ---------------------------------------------------------------------------
# Patch: output_dir=None means same folder as each input
# ---------------------------------------------------------------------------

_orig_convert_pdf = convert_pdf


def convert_pdf(  # noqa: F811
    pdf_path: Path,
    output_dir: Path | None,
    **kwargs,
) -> dict:
    effective_dir = output_dir if output_dir is not None else pdf_path.parent
    effective_dir.mkdir(parents=True, exist_ok=True)
    return _orig_convert_pdf(pdf_path=pdf_path, output_dir=effective_dir, **kwargs)


if __name__ == "__main__":
    main()
